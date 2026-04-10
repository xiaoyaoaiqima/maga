"""
文档解析服务 - 支持从 PDF/Word/Excel/PPT 中提取文本内容
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from loguru import logger

try:
    import openpyxl
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False
    logger.warning("openpyxl 未安装，Excel 解析功能不可用")

try:
    from PyPDF2 import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PdfReader = None  # type: ignore
    PDF_SUPPORT = False
    logger.warning("PyPDF2 未安装，PDF 解析功能不可用")

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    Document = None  # type: ignore
    DOCX_SUPPORT = False
    logger.warning("python-docx 未安装，Word 解析功能不可用")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    Presentation = None  # type: ignore
    PPTX_SUPPORT = False
    logger.warning("python-pptx 未安装，PPT 解析功能不可用")


class DocumentParser:
    """文档解析器"""

    @staticmethod
    def get_supported_types() -> list[str]:
        """获取支持的文件类型"""
        types = []
        if PDF_SUPPORT:
            types.extend(["pdf"])
        if DOCX_SUPPORT:
            types.extend(["doc", "docx"])
        if PPTX_SUPPORT:
            types.extend(["ppt", "pptx"])
        if EXCEL_SUPPORT:
            types.extend(["xls", "xlsx"])
        return types

    @staticmethod
    def is_supported(file_type: str) -> bool:
        """检查文件类型是否支持"""
        return file_type.lower() in DocumentParser.get_supported_types()

    @staticmethod
    async def parse_file(
        file_path: str,
        file_type: Optional[str] = None,
    ) -> dict:
        """
        解析文档文件，提取文本内容

        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，自动从扩展名推断）

        Returns:
            解析结果：{
                "success": bool,
                "text": str,
                "pages": int,  # PDF 页数
                "error": str | None
            }
        """
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "text": "",
                "error": f"文件不存在: {file_path}",
            }

        # 自动推断文件类型
        if file_type is None:
            ext = path.suffix.lstrip(".").lower()
            file_type = ext if ext else "unknown"

        file_type = file_type.lower()

        try:
            if file_type == "pdf" and PDF_SUPPORT:
                return DocumentParser._parse_pdf(file_path)
            elif file_type in ["doc", "docx"] and DOCX_SUPPORT:
                return DocumentParser._parse_docx(file_path)
            elif file_type in ["ppt", "pptx"] and PPTX_SUPPORT:
                return DocumentParser._parse_pptx(file_path)
            elif file_type in ["xls", "xlsx"] and EXCEL_SUPPORT:
                return DocumentParser._parse_excel(file_path)
            else:
                supported = ", ".join(DocumentParser.get_supported_types())
                return {
                    "success": False,
                    "text": "",
                    "error": f"不支持的文件类型: {file_type}，支持: {supported}",
                }
        except Exception as e:
            logger.error(f"文档解析失败: {e}")
            return {
                "success": False,
                "text": "",
                "error": f"解析失败: {str(e)}",
            }

    @staticmethod
    def _parse_pdf(file_path: str) -> dict:
        """解析 PDF 文件"""
        reader = PdfReader(file_path)
        pages = len(reader.pages)
        text_parts = []

        for page in reader.pages:
            try:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text)
            except Exception:
                continue

        text = "\n\n".join(text_parts)

        return {
            "success": True,
            "text": text,
            "pages": pages,
            "error": None,
        }

    @staticmethod
    def _parse_docx(file_path: str) -> dict:
        """解析 Word 文件"""
        doc = Document(file_path)
        text_parts = []

        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # 提取表格文本
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        text = "\n".join(text_parts)

        return {
            "success": True,
            "text": text,
            "pages": 1,
            "error": None,
        }

    @staticmethod
    def _parse_pptx(file_path: str) -> dict:
        """解析 PowerPoint 文件"""
        prs = Presentation(file_path)
        text_parts = []

        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

            if slide_texts:
                text_parts.append("\n".join(slide_texts))

        text = "\n\n".join(text_parts)

        return {
            "success": True,
            "text": text,
            "pages": len(prs.slides),
            "error": None,
        }

    @staticmethod
    def _parse_excel(file_path: str) -> dict:
        """解析 Excel 文件"""
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_texts = []

            for row in sheet.iter_rows(values_only=True):
                # 过滤空行
                row_values = [str(cell).strip() if cell is not None else "" for cell in row]
                if any(row_values):
                    sheet_texts.append(" | ".join(row_values))

            if sheet_texts:
                text_parts.append(f"=== {sheet_name} ===\n" + "\n".join(sheet_texts))

        text = "\n\n".join(text_parts)

        return {
            "success": True,
            "text": text,
            "pages": 1,
            "error": None,
        }
