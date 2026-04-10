"""
文档解析服务 - 支持 PDF、Word、PPT、Excel、CSV 等多种格式

功能：
1. 提取文档文本内容
2. 支持批量解析
3. 异步处理
4. 支持从 COS 下载文件后解析
"""
from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path
from typing import Optional

import openpyxl
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.knowledge_base_file import KnowledgeBaseFile
from app.services.excel_parser import ExcelParser
from app.services.cos_service import get_cos_service

logger = logging.getLogger(__name__)

# 临时目录配置
TEMP_DIR = Path(tempfile.gettempdir()) / "raap_knowledge_base_files"


async def _get_local_file_path(knowledge_base_file: KnowledgeBaseFile) -> Path:
    """
    获取文件本地路径

    如果 file_path 是 COS URL，则下载到临时目录
    如果是本地路径，直接返回

    Args:
        knowledge_base_file: 知识库文件记录

    Returns:
        本地文件路径
    """
    file_path_str = knowledge_base_file.file_path

    # 判断是否为 COS URL（以 http:// 或 https:// 开头）
    if file_path_str.startswith(("http://", "https://")):
        # COS URL，需要下载到临时目录
        cos_service = get_cos_service()
        local_path = TEMP_DIR / f"{knowledge_base_file.id}_{knowledge_base_file.file_name}"
        TEMP_DIR.mkdir(parents=True, exist_ok=True)

        success = await cos_service.download_file(file_path_str, local_path)
        if not success:
            raise FileNotFoundError(f"无法从 COS 下载文件: {file_path_str}")

        logger.info(f"COS 文件已下载到: {local_path}")
        return local_path
    else:
        # 本地路径，直接返回
        return Path(file_path_str)


async def _cleanup_temp_file(file_path: Path) -> None:
    """
    清理临时文件

    Args:
        file_path: 要清理的文件路径
    """
    try:
        if file_path.exists() and str(file_path).startswith(str(TEMP_DIR)):
            file_path.unlink()
            logger.info(f"临时文件已清理: {file_path}")
    except Exception as e:
        logger.warning(f"清理临时文件失败: {e}")


class DocumentParser:
    """文档解析器 - 支持多种格式"""

    # 支持的文件类型映射
    TYPE_MAPPING = {
        "pdf": "pdf",
        "doc": "word",
        "docx": "word",
        "ppt": "ppt",
        "pptx": "ppt",
        "xls": "excel",
        "xlsx": "excel",
        "csv": "csv",
    }

    @classmethod
    def get_doc_type(cls, file_ext: str) -> str:
        """获取文档类型"""
        return cls.TYPE_MAPPING.get(file_ext.lower(), "unknown")

    @classmethod
    def is_supported(cls, file_ext: str) -> bool:
        """检查文件类型是否支持"""
        return file_ext.lower() in cls.TYPE_MAPPING

    @classmethod
    async def extract_text(
        cls, file_path: str, file_type: str
    ) -> tuple[Optional[str], Optional[str]]:
        """
        提取文档文本内容

        Args:
            file_path: 文件路径
            file_type: 文件类型（pdf/word/ppt/excel/csv）

        Returns:
            (文本内容, 错误信息)
        """
        if not Path(file_path).exists():
            return None, f"文件不存在: {file_path}"

        try:
            if file_type == "excel":
                return await cls._extract_from_excel(file_path)
            elif file_type == "csv":
                return await cls._extract_from_csv(file_path)
            elif file_type == "pdf":
                return await cls._extract_from_pdf(file_path)
            elif file_type == "word":
                return await cls._extract_from_word(file_path)
            elif file_type == "ppt":
                return await cls._extract_from_ppt(file_path)
            else:
                return None, f"不支持的文件类型: {file_type}"
        except Exception as e:
            logger.error(f"解析文件失败 {file_path}: {e}")
            return None, str(e)

    @classmethod
    async def _extract_from_excel(cls, file_path: str) -> tuple[Optional[str], None]:
        """从 Excel 提取文本"""
        workbook = openpyxl.load_workbook(file_path, read_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(
                    str(cell.value) if cell.value is not None else ""
                    for cell in row
                ).strip()
                if row_text:
                    text_parts.append(row_text)

        workbook.close()
        return "\n".join(text_parts), None

    @classmethod
    async def _extract_from_csv(cls, file_path: str) -> tuple[Optional[str], None]:
        """从 CSV 提取文本"""
        import csv

        text_parts = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " ".join(row).strip()
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts), None

    @classmethod
    async def _extract_from_pdf(cls, file_path: str) -> tuple[Optional[str], Optional[str]]:
        """从 PDF 提取文本"""
        try:
            import PyPDF2

            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text.strip())

            return "\n".join(text_parts), None
        except ImportError:
            return None, "PyPDF2 未安装，无法解析 PDF"
        except Exception as e:
            return None, f"PDF 解析失败: {str(e)}"

    @classmethod
    async def _extract_from_word(cls, file_path: str) -> tuple[Optional[str], Optional[str]]:
        """从 Word 提取文本"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            for paragraph in doc.paragraphs:
                if paragraph.text and paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = " ".join(cell.text for cell in row.cells).strip()
                    if row_text:
                        text_parts.append(row_text)

            return "\n".join(text_parts), None
        except ImportError:
            return None, "python-docx 未安装，无法解析 Word"
        except Exception as e:
            return None, f"Word 解析失败: {str(e)}"

    @classmethod
    async def _extract_from_ppt(cls, file_path: str) -> tuple[Optional[str], Optional[str]]:
        """从 PPT 提取文本"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)

            return "\n".join(text_parts), None
        except ImportError:
            return None, "python-pptx 未安装，无法解析 PPT"
        except Exception as e:
            return None, f"PPT 解析失败: {str(e)}"


async def parse_knowledge_base_file(
    knowledge_base_file: KnowledgeBaseFile,
    db: AsyncSession,
) -> dict:
    """
    解析知识库文件中的文件

    Args:
        knowledge_base_file: 知识库文件记录
        db: 数据库会话

    Returns:
        解析结果 {success: bool, text_length: int, error?: string}
    """
    file_ext = Path(knowledge_base_file.file_name).suffix.replace(".", "").lower()
    doc_type = DocumentParser.get_doc_type(file_ext)

    # 获取本地文件路径（如果是 COS URL 会先下载）
    local_file_path = await _get_local_file_path(knowledge_base_file)
    is_temp_file = str(local_file_path).startswith(str(TEMP_DIR))

    try:
        text, error = await DocumentParser.extract_text(
            str(local_file_path), doc_type
        )

        if error:
            return {"success": False, "error": error}

        # 更新知识库文件状态
        from app.services.file_pool_service import KnowledgeBaseFileService

        service = KnowledgeBaseFileService(db)
        await service.update_parse_result(
            knowledge_base_file.id,
            parsed_count=len(text) if text else 0,
            total_count=len(text) if text else 0,
        )

        return {
            "success": True,
            "text_length": len(text) if text else 0,
            "file_type": doc_type,
        }
    finally:
        # 清理临时文件
        if is_temp_file:
            await _cleanup_temp_file(local_file_path)
