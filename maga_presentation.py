#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MAGA AI营销内容生成平台 - 汇报PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================
# 配色方案 - 深色科技风
# ============================================================
BG_DARK = RGBColor(0x12, 0x12, 0x1E)       # 深蓝黑背景 #12121E
BG_CARD = RGBColor(0x1E, 0x1E, 0x2E)       # 卡片背景 #1E1E2E
BG_CARD_ALT = RGBColor(0x25, 0x25, 0x3A)   # 交替卡片 #25253A
ACCENT = RGBColor(0xF5, 0xA6, 0x23)        # 金色强调 #F5A623
ACCENT_RED = RGBColor(0xE9, 0x45, 0x60)    # 红色强调 #E94560
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # 蓝色 #3B82F6
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)  # 绿色 #10B981
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # 纯白
TEXT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)     # 浅灰
TEXT_MUTED = RGBColor(0x6B, 0x6B, 0x7B)    # 暗灰
LINE_COLOR = RGBColor(0x3A, 0x3A, 0x4A)    # 分割线

# ============================================================
# 尺寸常量 (16:9, 13.333 x 7.5 inches)
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.6)
MARGIN_RIGHT = Inches(0.6)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)

# ============================================================
# 辅助函数
# ============================================================
def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=TEXT_WHITE,
                bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei', line_spacing=1.5):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.line_spacing = line_spacing
    # 设置中文字体
    run = p.runs[0]
    run.font._element.set('{http://schemas.openxmlformats.org/drawingml/2006/main}altLang', 'zh-CN')
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16, font_color=TEXT_GRAY,
                       bold=False, align=PP_ALIGN.LEFT, line_spacing=1.6):
    """添加多行文本（每段一个paragraph）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = align
        p.line_spacing = line_spacing
    return txBox

def add_bullet_text(slide, left, top, width, height, bullets, font_size=14, font_color=TEXT_GRAY,
                    bullet_color=ACCENT, line_spacing=1.5):
    """添加带项目符号的文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = 'Microsoft YaHei'
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
    return txBox

def add_line(slide, x1, y1, x2, y2, color=LINE_COLOR, width=Pt(1)):
    """添加线条"""
    line = slide.shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    return line

# ============================================================
# 创建演示文稿
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# 使用空白布局
blank_layout = prs.slide_layouts[6]  # blank

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BG_DARK)

# 左侧装饰条
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT)

# 顶部细线
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(2.5), Pt(2), ACCENT)

# 主标题
add_textbox(slide1, Inches(0.6), Inches(2.3), Inches(10), Inches(1.2),
            "MAGA", font_size=72, font_color=TEXT_WHITE, bold=True)

# 副标题
add_textbox(slide1, Inches(0.6), Inches(3.5), Inches(10), Inches(0.8),
            "AI 营销内容生成平台", font_size=40, font_color=ACCENT, bold=True)

# 描述
add_textbox(slide1, Inches(0.6), Inches(4.4), Inches(10), Inches(0.6),
            "让小红书内容生产从「人驱动」变为「AI 驱动」",
            font_size=20, font_color=TEXT_GRAY)

# 底部信息条
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.7), BG_CARD)
add_textbox(slide1, Inches(0.6), Inches(6.9), Inches(8), Inches(0.5),
            "技术团队 · 2026年6月", font_size=14, font_color=TEXT_MUTED)

# 右侧装饰圆
add_shape(slide1, MSO_SHAPE.OVAL, Inches(10.5), Inches(1.5), Inches(3), Inches(3),
          fill_color=None, line_color=ACCENT, line_width=Pt(2))
add_shape(slide1, MSO_SHAPE.OVAL, Inches(11.0), Inches(2.0), Inches(2), Inches(2),
          fill_color=None, line_color=ACCENT_BLUE, line_width=Pt(1))

# ============================================================
# 第2页：痛点 & 价值
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BG_DARK)

# 顶部标题区
add_textbox(slide2, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "01  痛点与价值", font_size=28, font_color=ACCENT, bold=True)
add_textbox(slide2, MARGIN_LEFT, Inches(0.9), Inches(8), Inches(0.5),
            "为什么需要 MAGA？", font_size=20, font_color=TEXT_WHITE, bold=True)

# 分隔线
add_shape(slide2, MSO_SHAPE.RECTANGLE, MARGIN_LEFT, Inches(1.35), Inches(12), Pt(1), LINE_COLOR)

# 左侧：传统方式
col1_x = MARGIN_LEFT
col_width = Inches(5.8)
add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, col1_x, Inches(1.6), col_width, Inches(5.4), BG_CARD)
add_textbox(slide2, col1_x + Inches(0.3), Inches(1.85), col_width - Inches(0.6), Inches(0.5),
            "传统方式", font_size=22, font_color=ACCENT_RED, bold=True)

problems = [
    "写手成本高，质量不稳定",
    "规则越写越重，维护困难",
    "风格不统一，品牌感弱",
    "反馈无法沉淀，经验流失"
]
for i, p in enumerate(problems):
    y = Inches(2.5 + i * 0.9)
    add_shape(slide2, MSO_SHAPE.RECTANGLE, col1_x + Inches(0.3), y, Pt(4), Pt(20), ACCENT_RED)
    add_textbox(slide2, col1_x + Inches(0.5), y - Pt(2), col_width - Inches(0.8), Inches(0.5),
                p, font_size=16, font_color=TEXT_GRAY)

# 右侧：MAGA 方案
col2_x = Inches(6.9)
add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, col2_x, Inches(1.6), col_width, Inches(5.4), BG_CARD)
add_textbox(slide2, col2_x + Inches(0.3), Inches(1.85), col_width - Inches(0.6), Inches(0.5),
            "MAGA 方案", font_size=22, font_color=ACCENT_GREEN, bold=True)

solutions = [
    "AI 一键生成，7×24 产出",
    "轻规则 + 多示例，运营好维护",
    "规则包驱动，风格可控",
    "反馈闭环，持续优化规则和示例"
]
for i, s in enumerate(solutions):
    y = Inches(2.5 + i * 0.9)
    add_shape(slide2, MSO_SHAPE.RECTANGLE, col2_x + Inches(0.3), y, Pt(4), Pt(20), ACCENT_GREEN)
    add_textbox(slide2, col2_x + Inches(0.5), y - Pt(2), col_width - Inches(0.8), Inches(0.5),
                s, font_size=16, font_color=TEXT_GRAY)

# 中间箭头
arrow = add_shape(slide2, MSO_SHAPE.RIGHT_ARROW, Inches(6.3), Inches(3.8), Inches(0.7), Inches(0.5), ACCENT)

# ============================================================
# 第3页：产品定位
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, BG_DARK)

# 顶部标题
add_textbox(slide3, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "02  产品定位", font_size=28, font_color=ACCENT, bold=True)

# 核心金句 - 大字号居中
quote_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.2), BG_CARD)

add_textbox(slide3, Inches(1.2), Inches(2.2), Inches(11), Inches(0.8),
            "帮品牌和博主用「最少输入」，", font_size=36, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide3, Inches(1.2), Inches(3.0), Inches(11), Inches(0.8),
            "稳定拿到「可直接发布」的小红书内容", font_size=36, font_color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# 三个核心价值点
cards = [
    ("极简前台", "仅需3个输入：产品/主题 + 目标人群 + 风格", ACCENT_BLUE),
    ("智能后台", "策略补全、风格优化、可发布性检查全自动", ACCENT_GREEN),
    ("直接可用", "输出 = 标题 + 正文，无需二次加工", ACCENT),
]
for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.8 + i * 4.2)
    add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.3), Inches(3.8), Inches(1.6), BG_CARD_ALT)
    add_shape(slide3, MSO_SHAPE.RECTANGLE, x, Inches(5.3), Inches(3.8), Pt(4), color)
    add_textbox(slide3, x + Inches(0.2), Inches(5.5), Inches(3.4), Inches(0.5),
                title, font_size=20, font_color=color, bold=True)
    add_textbox(slide3, x + Inches(0.2), Inches(6.0), Inches(3.4), Inches(0.8),
                desc, font_size=14, font_color=TEXT_GRAY)

# ============================================================
# 第4页：核心功能全景
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, BG_DARK)

add_textbox(slide4, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "03  核心功能全景", font_size=28, font_color=ACCENT, bold=True)

# 三栏结构
layers = [
    ("前台：极简任务页", "用户仅需输入产品/主题、目标人群、风格，一键生成", ACCENT_BLUE, [
        "单篇生成", "批量生成", "结果复制"
    ]),
    ("后台：智能生成引擎", "意图理解 → 策略补全 → 内容生成 → 风格优化 → 审核改写", ACCENT_GREEN, [
        "业务规则包驱动", "系统关键词自动补齐", "Expert 配置管理"
    ]),
    ("闭环：资产 & 质量", "业务规则 · 系统关键词 · 违禁词 · 反馈记录", ACCENT, [
        "违禁词审核", "人工反馈", "质量看板"
    ]),
]

for i, (title, desc, color, items) in enumerate(layers):
    x = Inches(0.5 + i * 4.3)
    # 卡片背景
    add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(4.0), Inches(5.8), BG_CARD)
    # 顶部色条
    add_shape(slide4, MSO_SHAPE.RECTANGLE, x, Inches(1.3), Inches(4.0), Pt(6), color)
    # 编号圆圈
    add_shape(slide4, MSO_SHAPE.OVAL, x + Inches(0.2), Inches(1.6), Inches(0.5), Inches(0.5), color)
    add_textbox(slide4, x + Inches(0.2), Inches(1.65), Inches(0.5), Inches(0.4),
                str(i+1), font_size=18, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide4, x + Inches(0.8), Inches(1.6), Inches(3.0), Inches(0.5),
                title, font_size=18, font_color=color, bold=True)
    # 描述
    add_textbox(slide4, x + Inches(0.2), Inches(2.2), Inches(3.6), Inches(0.8),
                desc, font_size=13, font_color=TEXT_GRAY)
    # 列表
    for j, item in enumerate(items):
        y = Inches(3.2 + j * 0.7)
        add_shape(slide4, MSO_SHAPE.RECTANGLE, x + Inches(0.3), y + Pt(6), Pt(3), Pt(3), color)
        add_textbox(slide4, x + Inches(0.5), y, Inches(3.3), Inches(0.5),
                    item, font_size=14, font_color=TEXT_GRAY)

# 连接箭头（下箭头）
for i in range(2):
    x = Inches(4.4 + i * 4.3)
    add_shape(slide4, MSO_SHAPE.DOWN_ARROW, x, Inches(3.5), Inches(0.5), Inches(0.5), ACCENT)

# ============================================================
# 第5页：技术架构
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, BG_DARK)

add_textbox(slide5, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "04  技术架构", font_size=28, font_color=ACCENT, bold=True)
add_textbox(slide5, MARGIN_LEFT, Inches(0.9), Inches(8), Inches(0.4),
            "MAGA 控制平面 + Hermes 执行平面", font_size=16, font_color=TEXT_GRAY)

# 三层架构图 - 用形状模拟
box_w = Inches(10.5)
box_h = Inches(1.4)
box_x = Inches(1.3)

# 第一层：控制台
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, box_x, Inches(1.5), box_w, box_h, BG_CARD_ALT)
add_textbox(slide5, box_x + Inches(0.3), Inches(1.7), Inches(3), Inches(0.4),
            "MAGA Console", font_size=18, font_color=ACCENT_BLUE, bold=True)
add_textbox(slide5, box_x + Inches(0.3), Inches(2.05), Inches(9), Inches(0.5),
            "Vue/Vite 前端 · 极简任务页 · 批量报告 · 资产管理 · 人工审核",
            font_size=14, font_color=TEXT_GRAY)

# 连接箭头
add_shape(slide5, MSO_SHAPE.DOWN_ARROW, box_x + box_w/2 - Inches(0.25), Inches(2.95), Inches(0.5), Inches(0.35), ACCENT)

# 第二层：后端
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, box_x, Inches(3.35), box_w, box_h, BG_CARD_ALT)
add_textbox(slide5, box_x + Inches(0.3), Inches(3.55), Inches(3), Inches(0.4),
            "MAGA Backend", font_size=18, font_color=ACCENT_GREEN, bold=True)
add_textbox(slide5, box_x + Inches(0.3), Inches(3.9), Inches(9), Inches(0.5),
            "FastAPI · MySQL · Redis · 调度器 + 状态机 · 资产中心 · Trace/审计系统",
            font_size=14, font_color=TEXT_GRAY)

# 连接箭头
add_shape(slide5, MSO_SHAPE.DOWN_ARROW, box_x + box_w/2 - Inches(0.25), Inches(4.8), Inches(0.5), Inches(0.35), ACCENT)

# 第三层：Worker
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, box_x, Inches(5.2), box_w, box_h, BG_CARD_ALT)
add_textbox(slide5, box_x + Inches(0.3), Inches(5.4), Inches(4), Inches(0.4),
            "Hermes maga-worker", font_size=18, font_color=ACCENT, bold=True)
add_textbox(slide5, box_x + Inches(0.3), Inches(5.75), Inches(9), Inches(0.5),
            "AI 执行引擎 · LLM 调用 · GE/AE 编排 · 内容生成 · 评分改写",
            font_size=14, font_color=TEXT_GRAY)

# 右侧职责说明
add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(1.5), Inches(0.9), Inches(5.1), BG_CARD)
add_textbox(slide5, Inches(0.2), Inches(3.0), Inches(0.8), Inches(2.5),
            "控\n制\n平\n面", font_size=16, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 第6页：业务闭环
# ============================================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, BG_DARK)

add_textbox(slide6, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "05  业务闭环", font_size=28, font_color=ACCENT, bold=True)
add_textbox(slide6, MARGIN_LEFT, Inches(0.9), Inches(8), Inches(0.4),
            "从生成到优化的完整飞轮", font_size=16, font_color=TEXT_GRAY)

# 四个步骤 - 横向流程
steps = [
    ("生成", "业务规则包\n+ 系统关键词\n+ Expert 配置", ACCENT_BLUE, "01"),
    ("审核", "违禁词扫描\n自动改写\n二次扫描兜底", ACCENT_GREEN, "02"),
    ("反馈", "通过 / 要求修改\n人工改写 / 批注\n沉淀训练素材", ACCENT, "03"),
    ("治理", "反馈聚合\n规则/示例调整\n质量看板跟踪", ACCENT_RED, "04"),
]

step_w = Inches(2.8)
step_h = Inches(3.5)
start_x = Inches(0.6)

for i, (title, desc, color, num) in enumerate(steps):
    x = start_x + i * (step_w + Inches(0.5))
    # 卡片
    add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), step_w, step_h, BG_CARD)
    # 顶部色条
    add_shape(slide6, MSO_SHAPE.RECTANGLE, x, Inches(1.6), step_w, Pt(5), color)
    # 编号
    add_shape(slide6, MSO_SHAPE.OVAL, x + Inches(0.15), Inches(1.85), Inches(0.5), Inches(0.5), color)
    add_textbox(slide6, x + Inches(0.15), Inches(1.9), Inches(0.5), Inches(0.4),
                num, font_size=16, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide6, x + Inches(0.75), Inches(1.88), Inches(1.8), Inches(0.4),
                title, font_size=20, font_color=color, bold=True)
    # 描述
    add_textbox(slide6, x + Inches(0.15), Inches(2.6), Inches(2.5), Inches(2.0),
                desc, font_size=14, font_color=TEXT_GRAY, align=PP_ALIGN.LEFT)

    # 箭头（除了最后一个）
    if i < len(steps) - 1:
        arrow_x = x + step_w + Inches(0.05)
        add_shape(slide6, MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(3.1), Inches(0.4), Inches(0.4), ACCENT)

# 底部飞轮说明
add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.4), Inches(12), Inches(1.0), BG_CARD_ALT)
add_textbox(slide6, Inches(1.0), Inches(5.6), Inches(11), Inches(0.7),
            "核心逻辑：生成质量 → 人工反馈 → 规则/示例调整 → 质量看板跟踪 → 更高质量的下一轮生成",
            font_size=16, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 第7页：关键能力
# ============================================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, BG_DARK)

add_textbox(slide7, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "06  关键能力", font_size=28, font_color=ACCENT, bold=True)

# 2x2 网格
capabilities = [
    ("业务规则包", "一种方向 = 一段轻规则 + 一组多示例", ACCENT_BLUE,
     ["规则只定方向、边界、不要做什么", "否定式规则比强制约束更有效", "示例教表达、供多样性、当语料", "整个语料块 = 给模型的提示词"]),
    ("系统提示词关键词", "版本化管理，自动从启用类别中选择子关键词", ACCENT_GREEN,
     ["人设 / 生文指令", "扰动规则 / 写作手法", "格式控制", "支持固定或轮换模式"]),
    ("Expert 配置", "提示词模板 + 模型参数的一次执行配置单元", ACCENT,
     ["article_generator_v1", "comment_generator_v1", "content_rewrite_v1", "模型/温度/Token 可调"]),
    ("违禁词审核改写", "确定性扫描 + AI 自然改写 + 二次扫描兜底", ACCENT_RED,
     ["系统违禁词库", "业务违禁词库", "命中后自动改写", "不通过则人工介入"]),
]

for i, (title, desc, color, items) in enumerate(capabilities):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.2 + row * 3.0)

    add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.0), Inches(2.7), BG_CARD)
    add_shape(slide7, MSO_SHAPE.RECTANGLE, x, y, Pt(5), Inches(2.7), color)

    add_textbox(slide7, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.4),
                title, font_size=18, font_color=color, bold=True)
    add_textbox(slide7, x + Inches(0.3), y + Inches(0.55), Inches(5.4), Inches(0.4),
                desc, font_size=13, font_color=TEXT_GRAY)

    for j, item in enumerate(items):
        item_y = y + Inches(1.0 + j * 0.4)
        add_shape(slide7, MSO_SHAPE.RECTANGLE, x + Inches(0.3), item_y + Pt(5), Pt(3), Pt(3), color)
        add_textbox(slide7, x + Inches(0.45), item_y, Inches(5.2), Inches(0.35),
                    item, font_size=13, font_color=TEXT_GRAY)

# ============================================================
# 第8页：当前进展
# ============================================================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide8, BG_DARK)

add_textbox(slide8, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "07  当前进展", font_size=28, font_color=ACCENT, bold=True)

# 时间线 - 三个阶段
phases = [
    ("Phase 1\n已完成", [
        "Clean Schema 完成",
        "统一生成链路跑通",
        "规则包批量生成",
        "历史能力大清理",
    ], ACCENT_GREEN, Inches(1.5)),
    ("Phase 2\n基础版已完成\n体验优化中", [
        "业务规则包管理",
        "系统关键词管理",
        "Expert + 模型配置",
        "违禁词审核改写",
    ], ACCENT, Inches(5.5)),
    ("Phase 3\n后续规划", [
        "改写质量与对比",
        "评价反馈工作台",
        "批次质量看板",
        "多活动规则包",
    ], ACCENT_BLUE, Inches(9.5)),
]

# 时间线主轴
add_shape(slide8, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.8), Inches(10), Pt(4), LINE_COLOR)

for title, items, color, x in phases:
    # 节点圆
    add_shape(slide8, MSO_SHAPE.OVAL, x - Inches(0.2), Inches(3.65), Inches(0.4), Inches(0.4), color)
    # 阶段标题
    add_textbox(slide8, x - Inches(1.15), Inches(2.15), Inches(2.3), Inches(1.0),
                title, font_size=15, font_color=color, bold=True, align=PP_ALIGN.CENTER)
    # 详情列表
    for j, item in enumerate(items):
        y = Inches(4.15 + j * 0.45)
        add_textbox(slide8, x - Inches(1.2), y, Inches(2.4), Inches(0.4),
                    "• " + item, font_size=12, font_color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# 底部说明
add_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8), BG_CARD_ALT)
add_textbox(slide8, Inches(1.0), Inches(6.15), Inches(11), Inches(0.5),
            "当前目标：从 Demo 跑通推进到运营可持续使用，重点打磨改写质量、反馈工作台和批次质量治理",
            font_size=14, font_color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# 第9页：核心竞争力 — 轻规则 + 多示例
# ============================================================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide9, BG_DARK)

add_textbox(slide9, MARGIN_LEFT, Inches(0.3), Inches(6), Inches(0.5),
            "08  核心竞争力", font_size=28, font_color=ACCENT, bold=True)
add_textbox(slide9, MARGIN_LEFT, Inches(0.75), Inches(8), Inches(0.4),
            "轻规则 + 多示例：让模型听得懂，让运营维护得了", font_size=16, font_color=TEXT_GRAY)

# === 顶部：走过的弯路 → 验证的解法 ===
# 左侧：弯路
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.2), Inches(3.8), Inches(2.2), BG_CARD)
add_textbox(slide9, Inches(0.6), Inches(1.3), Inches(3.4), Inches(0.35),
            "❌ 走过的弯路", font_size=16, font_color=ACCENT_RED, bold=True)
wrong_pts = [
    "规则越写越重",
    "约束层层堆叠",
    "维护成本爆炸",
    "模型学得死板",
]
for j, txt in enumerate(wrong_pts):
    add_textbox(slide9, Inches(0.6), Inches(1.7 + j * 0.4), Inches(3.4), Inches(0.35),
                "• " + txt, font_size=12, font_color=TEXT_GRAY)

# 中间箭头
add_shape(slide9, MSO_SHAPE.RIGHT_ARROW, Inches(4.3), Inches(2.0), Inches(0.5), Inches(0.4), ACCENT)

# 中间：核心解法
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(1.2), Inches(3.5), Inches(2.2), BG_CARD_ALT)
add_textbox(slide9, Inches(5.1), Inches(1.3), Inches(3.1), Inches(0.35),
            "核心解法", font_size=16, font_color=ACCENT, bold=True)
add_textbox(slide9, Inches(5.1), Inches(1.7), Inches(3.1), Inches(0.35),
            "一种方向", font_size=20, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide9, Inches(5.1), Inches(2.05), Inches(3.1), Inches(0.35),
            "一段轻规则", font_size=20, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide9, Inches(5.1), Inches(2.4), Inches(3.1), Inches(0.35),
            "一组多示例", font_size=20, font_color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 右侧箭头
add_shape(slide9, MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(2.0), Inches(0.5), Inches(0.4), ACCENT)

# 右侧：结果
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.2), Inches(3.8), Inches(2.2), BG_CARD)
add_textbox(slide9, Inches(9.3), Inches(1.3), Inches(3.4), Inches(0.35),
            "✓ 验证的结果", font_size=16, font_color=ACCENT_GREEN, bold=True)
right_pts = [
    "运营问题在哪改哪",
    "模型从示例学模式",
    "生成自然不模板",
    "风格可控可预期",
]
for j, txt in enumerate(right_pts):
    add_textbox(slide9, Inches(9.3), Inches(1.7 + j * 0.4), Inches(3.4), Inches(0.35),
                "• " + txt, font_size=12, font_color=TEXT_GRAY)

# === 中间下半：轻规则 vs 多示例 详解 ===
# 左侧：轻规则
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.6), Inches(6.0), Inches(2.8), BG_CARD)
add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(3.6), Inches(6.0), Pt(4), ACCENT_BLUE)
add_textbox(slide9, Inches(0.6), Inches(3.7), Inches(5.6), Inches(0.35),
            "轻规则：只定三件事", font_size=16, font_color=ACCENT_BLUE, bold=True)

rule_items = [
    ("方向", "这个业务规则在聊什么"),
    ("边界", "别写成什么、别超出什么阶段"),
    ("否定式", "\"不要写成...\"比\"必须包含...\"更有效"),
]
for j, (label, desc) in enumerate(rule_items):
    y = Inches(4.15 + j * 0.55)
    add_shape(slide9, MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.35), Inches(0.35), ACCENT_BLUE)
    add_textbox(slide9, Inches(0.6), Inches(4.2 + j * 0.55), Inches(0.35), Inches(0.25),
                str(j+1), font_size=12, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide9, Inches(1.05), Inches(4.15 + j * 0.55), Inches(1.2), Inches(0.3),
                label, font_size=13, font_color=ACCENT_BLUE, bold=True)
    add_textbox(slide9, Inches(2.2), Inches(4.15 + j * 0.55), Inches(4.0), Inches(0.3),
                desc, font_size=13, font_color=TEXT_GRAY)

# 右侧：多示例
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(3.6), Inches(6.2), Inches(2.8), BG_CARD)
add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(3.6), Inches(6.2), Pt(4), ACCENT_GREEN)
add_textbox(slide9, Inches(6.9), Inches(3.7), Inches(5.8), Inches(0.35),
            "多示例：一组示例承担三个角色", font_size=16, font_color=ACCENT_GREEN, bold=True)

example_items = [
    ("① 教自然表达", "模型从真实评论/素材中学口语化写法"),
    ("② 供多样性", "不同角度、不同场景，避免千篇一律"),
    ("③ 当语义语料", "整个语料块 = 给模型的完整提示词"),
]
for j, (label, desc) in enumerate(example_items):
    add_textbox(slide9, Inches(6.9), Inches(4.15 + j * 0.55), Inches(1.8), Inches(0.3),
                label, font_size=13, font_color=ACCENT_GREEN, bold=True)
    add_textbox(slide9, Inches(8.6), Inches(4.15 + j * 0.55), Inches(4.0), Inches(0.3),
                desc, font_size=13, font_color=TEXT_GRAY)

# === 底部：运营工作流 ===
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.55), Inches(12.3), Inches(0.75), BG_CARD_ALT)
add_textbox(slide9, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
            "运营的工作流：结果不好 → 分析问题在哪 → 规则问题改规则 / 示例问题改示例 / 多样性问题加示例",
            font_size=14, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 第10页：总结 & 下一步
# ============================================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide10, BG_DARK)

# 左侧装饰条
add_shape(slide10, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT)

add_textbox(slide10, MARGIN_LEFT, Inches(0.4), Inches(6), Inches(0.6),
            "09  总结与下一步", font_size=28, font_color=ACCENT, bold=True)

# 总结
add_textbox(slide10, MARGIN_LEFT, Inches(1.1), Inches(12), Inches(0.5),
            "MAGA = 营销内容生成工作台 + 内容资产中心 + Agent 控制平面 + 质量评估系统",
            font_size=18, font_color=TEXT_WHITE, bold=True)

# 下一步行动项
actions = [
    ("Q2 重点", "打磨改写质量和修改对比，让运营能判断系统改得好不好", ACCENT_BLUE),
    ("Q3 目标", "完善评价反馈工作台、批次质量看板和多活动规则包扩展", ACCENT_GREEN),
    ("长期愿景", "成为营销内容生成领域最稳定、最可控、最易运营维护的业务工作台", ACCENT),
]

for i, (title, desc, color) in enumerate(actions):
    y = Inches(1.9 + i * 1.3)
    # 左侧色块
    add_shape(slide10, MSO_SHAPE.RECTANGLE, MARGIN_LEFT, y, Pt(6), Inches(1.0), color)
    # 内容背景
    add_shape(slide10, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT + Pt(6), y, Inches(11.5), Inches(1.0), BG_CARD)
    add_textbox(slide10, MARGIN_LEFT + Inches(0.3), y + Inches(0.1), Inches(2.5), Inches(0.4),
                title, font_size=16, font_color=color, bold=True)
    add_textbox(slide10, MARGIN_LEFT + Inches(0.3), y + Inches(0.45), Inches(10.5), Inches(0.5),
                desc, font_size=14, font_color=TEXT_GRAY)

# 底部 CTA
add_shape(slide10, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), SLIDE_WIDTH, Inches(1.0), BG_CARD_ALT)
add_textbox(slide10, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
            "感谢聆听 · 期待领导指导与支持", font_size=20, font_color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
output_path = "/Users/luxifa/maga/MAGA_AI营销内容生成平台_汇报.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页")
